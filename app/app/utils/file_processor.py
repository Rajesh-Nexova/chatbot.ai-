"""
File processor for extracting text from various file formats.
Supports PDF, Word documents, Excel spreadsheets, PowerPoint presentations, and plain text files.
"""

import os
from typing import Optional, Tuple
from pathlib import Path
import tempfile
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from app.utils.logger import logger


class FileProcessor:
    """Processes different file formats and extracts text content."""

    # Supported file extensions and their MIME types
    SUPPORTED_FORMATS = {
        # Text files
        '.txt': 'text/plain',
        '.md': 'text/markdown',
        '.csv': 'text/csv',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.html': 'text/html',
        '.htm': 'text/html',

        # Documents
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',

        # Spreadsheets
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',

        # Presentations
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
    }

    @classmethod
    def is_supported_format(cls, filename: str) -> bool:
        """Check if the file format is supported."""
        file_ext = Path(filename).suffix.lower()
        return file_ext in cls.SUPPORTED_FORMATS

    @classmethod
    def get_mime_type(cls, filename: str) -> str:
        """Get the MIME type for a file based on its extension."""
        file_ext = Path(filename).suffix.lower()
        return cls.SUPPORTED_FORMATS.get(file_ext, 'application/octet-stream')

    @classmethod
    async def extract_text(cls, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """
        Extract text from file content based on file type.

        Returns:
            Tuple of (extracted_text, is_text_extractable)
        """
        file_ext = Path(filename).suffix.lower()

        try:
            if file_ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm']:
                # Plain text files
                try:
                    text = file_content.decode('utf-8')
                    return text, True
                except UnicodeDecodeError:
                    return "", False

            elif file_ext == '.pdf':
                return cls._extract_pdf_text(file_content), True

            elif file_ext in ['.docx', '.doc']:
                return cls._extract_docx_text(file_content), True

            elif file_ext in ['.xlsx', '.xls']:
                return cls._extract_excel_text(file_content), True

            elif file_ext in ['.pptx', '.ppt']:
                return cls._extract_pptx_text(file_content), True

            else:
                # Unsupported format
                logger.warning(f"Unsupported file format: {filename}")
                return "", False

        except Exception as exc:
            logger.error(f"Error extracting text from {filename}: {exc}")
            return "", False

    @classmethod
    def _extract_pdf_text(cls, file_content: bytes) -> str:
        """Extract text from PDF file."""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(file_content)
            temp_file_path = temp_file.name

        try:
            reader = PdfReader(temp_file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        finally:
            os.unlink(temp_file_path)

    @classmethod
    def _extract_docx_text(cls, file_content: bytes) -> str:
        """Extract text from Word document."""
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
            temp_file.write(file_content)
            temp_file_path = temp_file.name

        try:
            doc = Document(temp_file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        finally:
            os.unlink(temp_file_path)

    @classmethod
    def _extract_excel_text(cls, file_content: bytes) -> str:
        """Extract text from Excel spreadsheet."""
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_file:
            temp_file.write(file_content)
            temp_file_path = temp_file.name

        try:
            wb = load_workbook(temp_file_path, data_only=True)
            text = ""

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"Sheet: {sheet_name}\n"

                for row in sheet.iter_rows(values_only=True):
                    # Convert all values to strings and join
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():  # Only add non-empty rows
                        text += row_text + "\n"

                text += "\n"

            return text.strip()
        finally:
            os.unlink(temp_file_path)

    @classmethod
    def _extract_pptx_text(cls, file_content: bytes) -> str:
        """Extract text from PowerPoint presentation."""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(file_content)
            temp_file_path = temp_file.name

        try:
            prs = Presentation(temp_file_path)
            text = ""

            for slide_number, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_number}:\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"

                text += "\n"

            return text.strip()
        finally:
            os.unlink(temp_file_path)